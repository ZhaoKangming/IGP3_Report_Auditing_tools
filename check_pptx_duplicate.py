# -*- coding:utf-8 -*-

import csv
import os
import shutil
import win32com.client as win32
from pptx import Presentation
import re


# 定义公共变量
script_path: str = os.path.dirname(os.path.realpath(__file__))
passed_report_folder: str = os.path.join(script_path, '..', 'reports','合格报告')
summary_csv_path: str = os.path.join(script_path, 'records', 'summary_data.csv')
summary_dict: dict = {}


def ppt_to_pptx():
    '''
    【功能】删除隐藏文件、转化ppt为pptx
    '''
    app = win32.gencache.EnsureDispatch('PowerPoint.Application')
    have_ppt: bool = False

    for file in os.listdir(passed_report_folder):
        file_path: str = os.path.join(passed_report_folder, file)
        if file[0] == '.':
            os.remove(file_path)
        else:
            if os.path.splitext(file)[1] == '.ppt':
                have_ppt = True
                print(f'正在转换文件 —— 《{file}》')
                new_pptx_path: str = file_path + 'x'
                backup_ppt_path: str = os.path.join(script_path, '..', 'reports', '合格ppt格式备份', file)
                shutil.copyfile(file_path, backup_ppt_path)   # 备份一下合格ppt报告
                office_obj = app.Presentations.Open(file_path, WithWindow=False)
                office_obj.SaveAs(new_pptx_path)
                office_obj.Close()
                os.remove(file_path)

    app.Quit()
    if have_ppt == True:
        print('转换文件格式完成！')


def load_summary_data():
    '''
    【功能】载入总结数据csv
    '''
    global summary_dict
    with open(summary_csv_path, 'r', encoding='utf-8-sig') as csv_reader:
        content = csv_reader.read()
        rows = content.split('\n')
        for row in rows:
            if row != '':
                summary_dict[row.split(',')[0]] = row.split(',')[1]
        summary_dict.pop('报告文件名')


def update_summary_data():
    '''
    【功能】更新报告的总结数据
    '''
    global summary_dict
    with open(summary_csv_path, 'a+', encoding='utf-8-sig', newline='') as csv_file:
        csv_writer = csv.writer(csv_file)
        for report_file in os.listdir(passed_report_folder):
            if report_file[0] != '.' and os.path.splitext(report_file)[1] == '.pptx':
                filename: str = os.path.splitext(report_file)[0]
                if not filename in summary_dict.keys():
                    content_dict: dict = get_pptx_content(os.path.join(passed_report_folder, report_file))
                    summary_result: list = get_content_summary(content_dict)
                    if summary_result[0] == True:
                        if summary_result[2] >= 190:  # 含标点字数大于等于190
                            summary_dict[filename] = summary_result[1]
                            csv_writer.writerow([filename, summary_result[1],summary_result[2]])
                        else:
                            print(f'[ERROR] --> 总结字数不足，含标点 {summary_result[2]} 字 ——《{report_file}》')
                    else:
                        print(f'[ERROR] --> 未找到获益与展望部分 ——《{report_file}》')


def get_pptx_content(pptx_path: str) -> dict:
    '''
    【功能】从报告中提取报告的总结部分
    【输出】内容字典，key为页码数，value为文本内容的连接
    :param pptx_path: 报告pptx的文件路径
    '''
    prs = Presentation(pptx_path)
    backup_content_dict: dict = {}
    content_dict: dict = {}
    i: int = 1

    for slide in prs.slides:
        text: str = ''
        backup_text: str = ''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    backup_text += run.text
                    text += run.text.replace(' ', '').replace('_', '')
        content_dict[i] = text
        backup_content_dict[i] = backup_text
        i += 1
    
    return content_dict


def get_content_summary(content_dict: dict) -> list:
    '''
    【功能】在报告内容字典中获取医生自己写的总结部分
    '''
    have_summary: bool = False
    have_typical_case: bool = False
    start_page_numb: int = 0        # 总结与展望开始的页码
    end_page_numb: int = 0          # 总结与展望结束的页码
    summary_text: str = ''          # 总结与展望的文本
    sentence_list: list = []
    character_numb: int = 0

    for i in range(5, len(content_dict)+1):     # 因为目录肯定出现在前五页，而获益与展望不会出现在前五页
        if '胰岛素规范实践的获益' in content_dict[i]:   # 因为有的医生会把获益与展望放到两页来写
            have_summary = True
            start_page_numb: int = i
            for j in range(i+1, i +4):
                if '典型病例分享' in content_dict[j]:
                    have_typical_case = True
                    end_page_numb: int = j - 1
                    break
            break
    
    if have_summary == True:
        if have_typical_case == False:
            end_page_numb: int = start_page_numb    # 如果在 “总结页” 后面的三页内都没有 “典型病例分享”，我们就默认总结只有一页
        # 合并总结部分的文本
        for i in range(start_page_numb, end_page_numb+1):
            summary_text += content_dict[i]

        # 清除提示性文本
        del_sentence_list: list = ['胰岛素规范实践的获益与展望',
                                    '胰岛素规范实践的获益',
                                    '胰岛素规范实践的展望',
                                    '请回顾本组患者依从性数据',
                                    '探讨规律随访的获益',
                                    '可从以下几个方面总结',
                                    '请回顾本组患者起始胰岛素治疗的特点',
                                    '探讨规范起始胰岛素的临床获益']
        for del_sentence in del_sentence_list:
            summary_text = summary_text.replace(del_sentence, '')

        character_numb = len(summary_text)  # 计算字数

        del_point_sentence_list: list = ['胰岛素起始时机',
                                        '口服药联用情况',
                                        '2型糖尿病病程',
                                        '并发症/合并症',
                                        '其他',
                                        '规范定期随访',
                                        '规律随访患者获益',
                                        '常见患者脱落原因',
                                        '糖化水平',
                                        '提升治疗依从性的经验']
        for del_point_sentence in del_point_sentence_list:
            summary_text = summary_text.replace(del_point_sentence, '')

        # 清除数字、英文字符、特殊标点
        pattern_1: str = '[a-zA-Z0-9#$%&()*+-/<=>@★、…【】[\\]^_`{|}~]+'
        summary_text = re.sub(pattern_1, '', summary_text)

        # 将语句拆分
        pattern_2 = r',|\.|/|;|\'|`|\[|\]|<|>|\?|:|：|"|\{|\}|\~|!|\(|\)|-|=|\_|\+|，|。|、|；|“|”|‘|’|·|！| |…|（|）'
        sentence_list = [i for i in re.split(pattern_2, summary_text) if i != '']       # 清除列表中的空值

    return [have_summary, sentence_list, character_numb]


def compute_similarity():
    '''
    【功能】计算总结部分的相似度
    '''
    pass


def initialized_database():
    '''
    【功能】初始化获益与展望数据库
    '''
    ppt_to_pptx()
    load_summary_data()
    update_summary_data()


#TODO:首先核对该医生的报告名称是否已经在字典中

initialized_database()
