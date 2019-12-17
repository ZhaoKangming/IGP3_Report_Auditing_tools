# -*- coding:UTF-8 -*-

from pptx import Presentation
import re
import os
from PyQt5.QtWidgets import *
from PyQt5.QtCore import Qt

#TODO:增加对ppt格式的支持
# 通用函数
def only_chinese(content):
    '''
    【功能】将传入的文本仅保留中文
    '''
    # 处理前进行相关的处理，包括转换成Unicode等
    pattern = re.compile('[^\u4e00-\u9fa50-9]')  # 中文的编码范围是：\u4e00到\u9fa5
    zh_str = "".join(pattern.split(content))
    return zh_str


# 【获取报告的文本内容】-----------------------------------------------------------------------------------

def get_pptx_content(pptx_path: str) -> dict:
    '''
    【功能】依据编号获取pptx文件内的文本内容
    :param pptx_path：临时报告pptx的路径
    '''
    prs = Presentation(pptx_path)
    backup_content_dict: dict = {}
    content_dict: dict = {}
    i: int = 0

    for slide in prs.slides:
        text: str = ''
        backup_text: str = ''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    backup_text += run.text
                    text += run.text.replace(' ', '').replace('_', '')
        content_dict[i] = text
        backup_content_dict[i] = backup_text
        i += 1

    print(content_dict)
    return content_dict

# 【审核报告的内容】-----------------------------------------------------------------------------------

def audit_slide(content_dict : dict, doc_name : str, pptx_path : str) -> dict:
    audit_result_dict: dict = {'审核结果':[], '修改记录':[], '错误记录':[]}
    slide_sort_dict: dict = {'首页': [], '注意事项': [], '内容目录': [], '基线情况汇总':[]}

    # ----------------------- 首页 -----------------------
    if '胰岛素规范临床实践' in content_dict[0]:  # 获取首页页码
        slide_sort_dict['首页'].append(0)
        temp_text: str = only_chinese(content_dict[0])

        # 检查汇报人姓名
        del_list: list = ['胰岛素规范临床实践', '总结', '报告', '汇报人']
        for del_str in del_list:
            temp_text = temp_text.replace(del_str, '')
        if len(temp_text) > 0:
            if doc_name not in temp_text:
                audit_result_dict['错误记录'].append('【首页】汇报人姓名与上传报告医生姓名不一致！')
    else:
        audit_result_dict['错误记录'].append("【首页】第一页没有'胰岛素规范临床实践'文本！")


    # ----------------------- 内容目录 -----------------------
    for i in [1,2]:
        if '注意事项' in content_dict[i]:  # 获取注意事项与内容目录页码
            slide_sort_dict['注意事项'].append(i)
        elif '治疗方案' in content_dict[i] or '病例分享' in content_dict[i]:
            slide_sort_dict['内容目录'].append(i)

    if slide_sort_dict['内容目录']:
        content_page_numb: int = slide_sort_dict['内容目录'][0]
        title_content_text: str = content_dict[content_page_numb]
        lack_title_list: list = []
        title_list: list = ['患者情况汇总', '治疗方案', '治疗结果', '典型病例分享', '胰岛素规范实践的获益', '胰岛素规范实践临床展望']
        for title_str in title_list:
            if not title_str in title_content_text:
                lack_title_list.append(title_str)
        if lack_title_list:
            lack_title_str: str = '、'.join(lack_title_list)
            audit_result_dict['错误记录'].append(f"【内容目录】缺少以下的模板中的字段：{lack_title_str}！")
    else:
        audit_result_dict['错误记录'].append('【内容目录】未发现内容目录页！')


    # ----------------------- 患者基线情况汇总 -----------------------
    # 获取患者基线情况目录页码
    for i in [2,3,4,5]:
        if '基线情况汇总' in content_dict[i]:  
            slide_sort_dict['基线情况汇总'].append(i)

    # 生成基线情况汇总字符串
    baseline_page_numbs: int = len(slide_sort_dict['基线情况汇总'])
    baseline_record: str = ''
    if baseline_page_numbs == 0 :
        audit_result_dict['错误记录'].append('【基线情况汇总】未发现标题含有-基线情况汇总-的页面！')
    else:
        for page_numb in slide_sort_dict['基线情况汇总']:
            baseline_record += content_dict[page_numb]

    # 检查基线情况汇总的填写正确与否
    baseline_record_no_punctuation: str = re.sub(r'[^\w\s]', '', baseline_record)
    # result = re.findall(".*纳入(.*)名患者男性(.*)名女性(.*)名平均", s1)
    # for x in result:
    #     print(x)

        
    # ----------------------- 注意事项 -----------------------
    if slide_sort_dict['注意事项']:
        # 删除'注意事项'页面
        audit_result_dict['修改记录'].append("【注意事项】删除了 注意事项 页面！")
        prs = Presentation(pptx_path)
        del_page_numb : int = slide_sort_dict['注意事项'][0]
        rId = prs.slides._sldIdLst[del_page_numb].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[del_page_numb]
        prs.save(pptx_path)

    return audit_result_dict
